"""Extractor de PowerPoint (.pptx)."""
from __future__ import annotations

from pathlib import Path

from ..models import ExtractedDoc
from .base import ExtractionError, Extractor


class PptxExtractor(Extractor):
    extensions = (".pptx",)

    def extract(self, path: Path) -> ExtractedDoc:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover
            raise ExtractionError("Falta dependencia: python-pptx") from exc

        prs = Presentation(str(path))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            lines.append(text)
            if lines:
                slides.append(f"# Diapositiva {i}\n" + "\n".join(lines))

        meta = {"format": "pptx", "slides": len(prs.slides)}
        return ExtractedDoc(text="\n\n".join(slides), meta=meta)
